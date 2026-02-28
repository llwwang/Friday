#!/usr/bin/env python3
"""
AI能力建设项目PPT生成器
为陛下制作专业的项目规划PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 添加空白幻灯片
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 设置深色背景
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)  # 深蓝黑色

# 添加标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "大数据运维AI能力建设项目"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(56, 189, 248)  # 亮蓝色
title_para.alignment = PP_ALIGN.CENTER

# 添加副标题
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Big Data Operations AI Capability Building"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(16)
subtitle_para.font.color.rgb = RGBColor(148, 163, 184)  # 灰色
subtitle_para.alignment = PP_ALIGN.CENTER

# 三个阶段的配置
phases = [
    {
        "title": "Phase 1",
        "subtitle": "for一线 - 工单机器人",
        "color": RGBColor(59, 130, 246),  # 蓝色
        "items": [
            "• 工单问答机器人",
            "• 基于现有知识库自动回应",
            "• 解决不了继续提单",
            "• 降低咨询工单量"
        ],
        "icon": "🤖"
    },
    {
        "title": "Phase 2",
        "subtitle": "for二线 - 内部问答平台",
        "color": RGBColor(34, 197, 94),  # 绿色
        "items": [
            "• 研发内部问答平台",
            "• 结合MCP协议",
            "• Agent技术拉通代码",
            "• 提升工单解决效率"
        ],
        "icon": "💻"
    },
    {
        "title": "Phase 3",
        "subtitle": "for客户 - 自动化运维分身",
        "color": RGBColor(168, 85, 247),  # 紫色
        "items": [
            "• OpenClaw + Ollama产品化",
            "• 大数据自动化运维分身",
            "• 私有化客户小模型",
            "• 客户内部运维自闭环"
        ],
        "icon": "🚀"
    }
]

# 绘制三个阶段卡片
card_width = Inches(4)
card_height = Inches(5)
start_x = Inches(0.5)
start_y = Inches(1.8)
gap = Inches(0.166)

for i, phase in enumerate(phases):
    x = start_x + i * (card_width + gap)
    
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, start_y, card_width, card_height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)  # 卡片背景色
    card.line.color.rgb = phase["color"]
    card.line.width = Pt(2)
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x, start_y, card_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = phase["color"]
    top_bar.line.fill.background()
    
    # Phase标题
    phase_title_box = slide.shapes.add_textbox(
        x + Inches(0.2), start_y + Inches(0.3), 
        card_width - Inches(0.4), Inches(0.5)
    )
    phase_title_frame = phase_title_box.text_frame
    phase_title_frame.text = phase["title"]
    phase_title_para = phase_title_frame.paragraphs[0]
    phase_title_para.font.size = Pt(28)
    phase_title_para.font.bold = True
    phase_title_para.font.color.rgb = phase["color"]
    
    # 副标题
    phase_subtitle_box = slide.shapes.add_textbox(
        x + Inches(0.2), start_y + Inches(0.8), 
        card_width - Inches(0.4), Inches(0.5)
    )
    phase_subtitle_frame = phase_subtitle_box.text_frame
    phase_subtitle_frame.text = phase["subtitle"]
    phase_subtitle_para = phase_subtitle_frame.paragraphs[0]
    phase_subtitle_para.font.size = Pt(16)
    phase_subtitle_para.font.bold = True
    phase_subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # 内容列表
    content_box = slide.shapes.add_textbox(
        x + Inches(0.2), start_y + Inches(1.4), 
        card_width - Inches(0.4), Inches(3)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for j, item in enumerate(phase["items"]):
        if j == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(203, 213, 225)  # 浅灰色
        p.space_after = Pt(12)

# 底部信息
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12.333), Inches(0.4))
footer_frame = footer_box.text_frame
footer_frame.text = "项目启动 | 2026年度重点 | 大数据运维团队"
footer_para = footer_frame.paragraphs[0]
footer_para.font.size = Pt(12)
footer_para.font.color.rgb = RGBColor(100, 116, 139)
footer_para.alignment = PP_ALIGN.CENTER

# 保存
output_path = "/tmp/AI能力建设项目规划.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 幻灯片尺寸: 16:9 宽屏")
print(f"🎨 主题风格: 深色科技风")
print(f"📋 内容: 3个Phase完整规划")
